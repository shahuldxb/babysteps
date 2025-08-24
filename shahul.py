#!/usr/bin/env python3
# ingest_docs_hardcoded.py — Incremental RAG ingest for mixed docs under C:\pp
# - Scans DOC/DOCX/PDF/CSV/TSV/XLS/XLSX/TXT/MD/JSON/PPT/PPTX/RTF
# - Extracts text, chunks, embeds, and upserts to PGVector
# - Skips files if content hash unchanged; re-ingests when changed
# - Keeps your exact keys/endpoints as provided

import os, sys, hashlib
from pathlib import Path
from typing import List, Dict
from datetime import datetime

from tqdm import tqdm
from docx import Document as DocxDocument

from langchain_openai import AzureOpenAIEmbeddings, AzureChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain_postgres import PGVector
from langchain.chains import create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_core.prompts import ChatPromptTemplate

# ======================
# HARD-CODED SETTINGS
# ======================
DOC_ROOT = Path(r"C:\pp")          # Folder of your documents
COLLECTION_NAME = "ppdocs"         # PGVector collection/table name

# ---- Azure OpenAI (KEEP EXACT VALUES) ----
AZURE_OPENAI_ENDPOINT = "https://shace-mejiegxb-eastus2.cognitiveservices.azure.com/"
AZURE_OPENAI_API_KEY  = "9LTf4DOFBLDVHcCm9LaEcJ1hLUK3QIhzDVqGOhKqh6nsF3VBIvoLJQQJ99BHACHYHv6XJ3w3AAAAACOG77Gk"

# Embeddings config
AZURE_OPENAI_API_VERSION = "2024-02-01"          # embeddings API version
AZURE_EMBED_DEPLOYMENT   = "text-embedding-3-large"

# Chat config
AZURE_CHAT_API_VERSION   = "2024-12-01-preview"  # chat API version
AZURE_CHAT_DEPLOYMENT    = "gpt-4o"

# --- Postgres (single string) ---
PG_CONN = "postgresql+psycopg://shahul:Cpu%4012345@shahulhannah.postgres.database.azure.com:5432/postgres?sslmode=require"

# ---- Chunking ----
CHUNK_SIZE = 2500
CHUNK_OVERLAP = 250
BATCH_SIZE = 64

# File types to ingest
SUPPORTED_EXTS = {
    ".docx", ".doc",
    ".pdf",
    ".csv", ".tsv",
    ".xlsx", ".xls",
    ".txt", ".md", ".json",
    ".pptx", ".ppt",
    ".rtf",
}
# ======================


def _is_missing(v: object) -> bool:
    if not isinstance(v, str):
        return True
    s = v.strip()
    return (s == "" or "PASTE" in s or "YOUR-" in s)


def extract_text_from_docx(path: Path) -> str:
    d = DocxDocument(path)
    parts: List[str] = []
    # paragraphs
    for p in d.paragraphs:
        txt = (p.text or "").strip()
        if txt:
            parts.append(txt)
    # tables
    for t in d.tables:
        for row in t.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            row_text = " | ".join([c for c in cells if c])
            if row_text:
                parts.append(row_text)
    return "\n".join(parts).strip()


def extract_text_from_pdf(path: Path) -> str:
    # Try pdfminer.six first
    try:
        from pdfminer.high_level import extract_text
        txt = extract_text(str(path)) or ""
        return txt.strip()
    except Exception as e:
        print(f"  pdfminer failed ({path.name}): {e}")
    # Fallback to PyPDF2
    try:
        import PyPDF2
        txt_parts = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    txt_parts.append(page.extract_text() or "")
                except Exception:
                    continue
        return "\n".join(txt_parts).strip()
    except Exception as e:
        print(f"  PyPDF2 failed ({path.name}): {e}")
    return ""


def extract_text_from_csv_tsv(path: Path) -> str:
    import csv
    sep = "\t" if path.suffix.lower() == ".tsv" else ","
    parts = []
    try:
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f, delimiter=sep)
            for row in reader:
                if any((cell or "").strip() for cell in row):
                    parts.append(" | ".join((cell or "").strip() for cell in row))
    except Exception as e:
        print(f"  CSV/TSV read failed ({path.name}): {e}")
    return "\n".join(parts).strip()


def extract_text_from_excel(path: Path) -> str:
    # Prefer pandas for multi-sheet handling
    try:
        import pandas as pd
        dfs = pd.read_excel(str(path), sheet_name=None, dtype=str)  # all sheets
        parts = []
        for sheet_name, df in dfs.items():
            if df is None or df.empty:
                continue
            df = df.fillna("")
            parts.append(f"## Sheet: {sheet_name}")
            for row in df.itertuples(index=False):
                row_vals = [str(v).strip() for v in row if str(v).strip()]
                if row_vals:
                    parts.append(" | ".join(row_vals))
        return "\n".join(parts).strip()
    except Exception as e:
        print(f"  pandas read_excel failed ({path.name}): {e}")
    # Fallback: openpyxl for .xlsx
    try:
        if path.suffix.lower() == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(filename=str(path), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"## Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if row_vals:
                        parts.append(" | ".join(row_vals))
            return "\n".join(parts).strip()
    except Exception as e:
        print(f"  openpyxl failed ({path.name}): {e}")
    # Fallback for legacy .xls (xlrd)
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets():
            parts.append(f"## Sheet: {sheet.name}")
            for r in range(sheet.nrows):
                row = [str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
                row = [v for v in row if v]
                if row:
                    parts.append(" | ".join(row))
        return "\n".join(parts).strip()
    except Exception as e:
        print(f"  xlrd failed ({path.name}): {e}")
    return ""


def extract_text_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts).strip()
    except Exception as e:
        print(f"  python-pptx failed ({path.name}): {e}")
    return ""


def extract_text_with_textract(path: Path) -> str:
    """Generic fallback for many formats (.doc, .ppt, sometimes .pdf) if textract is installed."""
    try:
        import textract
        data = textract.process(str(path))
        return data.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        print(f"  textract failed ({path.name}): {e}")
    return ""


def extract_text_from_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return rtf_to_text(f.read()).strip()
    except Exception as e:
        print(f"  RTF parse failed ({path.name}): {e}")
    return ""


def extract_text_plain(path: Path) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        print(f"  Plain text read failed ({path.name}): {e}")
    return ""


def extract_text_from_json(path: Path) -> str:
    try:
        import json
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            obj = json.load(f)
        # Simple flatten
        def _flatten(o):
            if isinstance(o, dict):
                for k, v in o.items():
                    yield f"{k}:"
                    yield from _flatten(v)
            elif isinstance(o, list):
                for i, v in enumerate(o):
                    yield f"- {str(v)[:200] + ('...' if len(str(v)) > 200 else '')}"
                    yield from _flatten(v)
            else:
                yield str(o)
        return "\n".join(x for x in _flatten(obj) if x).strip()
    except Exception as e:
        print(f"  JSON parse failed ({path.name}): {e}")
    return ""


def extract_text_from_any(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext in (".csv", ".tsv"):
        return extract_text_from_csv_tsv(path)
    if ext in (".xlsx", ".xls"):
        return extract_text_from_excel(path)
    if ext in (".txt", ".md"):
        return extract_text_plain(path)
    if ext == ".json":
        return extract_text_from_json(path)
    if ext == ".pptx":
        return extract_text_from_pptx(path)
    if ext in (".doc", ".ppt"):
        # Legacy Office formats: try textract if available
        return extract_text_with_textract(path)
    if ext == ".rtf":
        return extract_text_from_rtf(path)
    # Last-chance generic try via textract
    return extract_text_with_textract(path)


def hash_text(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8", errors="ignore")).hexdigest()


def build_embeddings_client() -> AzureOpenAIEmbeddings:
    return AzureOpenAIEmbeddings(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        openai_api_version=AZURE_OPENAI_API_VERSION,
        azure_deployment=AZURE_EMBED_DEPLOYMENT,
        model=AZURE_EMBED_DEPLOYMENT
    )


def build_chat_llm() -> AzureChatOpenAI:
    return AzureChatOpenAI(
        azure_endpoint=AZURE_OPENAI_ENDPOINT,
        api_key=AZURE_OPENAI_API_KEY,
        openai_api_version=AZURE_CHAT_API_VERSION,
        azure_deployment=AZURE_CHAT_DEPLOYMENT,
        model=AZURE_CHAT_DEPLOYMENT,
    )


def ensure_registry_table():
    """A tiny registry to skip unchanged files."""
    import sqlalchemy as sa
    from sqlalchemy import text
    engine = sa.create_engine(PG_CONN, pool_pre_ping=True)
    with engine.begin() as conn:
        conn.execute(text("""
        CREATE TABLE IF NOT EXISTS ingest_registry (
            source_path TEXT PRIMARY KEY,
            content_hash TEXT NOT NULL,
            file_mtime BIGINT NOT NULL,
            num_chunks INT NOT NULL,
            last_ingested_at TIMESTAMPTZ NOT NULL DEFAULT NOW()
        );
        """))


def get_existing_hash(source_path: str) -> str:
    import sqlalchemy as sa
    from sqlalchemy import text
    engine = sa.create_engine(PG_CONN, pool_pre_ping=True)
    with engine.begin() as conn:
        row = conn.execute(text(
            "SELECT content_hash FROM ingest_registry WHERE source_path = :p"
        ), {"p": source_path}).fetchone()
        return row[0] if row else ""


def upsert_registry(source_path: str, content_hash: str, file_mtime: int, num_chunks: int):
    import sqlalchemy as sa
    from sqlalchemy import text
    engine = sa.create_engine(PG_CONN, pool_pre_ping=True)
    with engine.begin() as conn:
        conn.execute(text("""
            INSERT INTO ingest_registry (source_path, content_hash, file_mtime, num_chunks, last_ingested_at)
            VALUES (:p, :h, :m, :n, NOW())
            ON CONFLICT (source_path) DO UPDATE
            SET content_hash = EXCLUDED.content_hash,
                file_mtime = EXCLUDED.file_mtime,
                num_chunks = EXCLUDED.num_chunks,
                last_ingested_at = NOW();
        """), {"p": source_path, "h": content_hash, "m": file_mtime, "n": num_chunks})


def list_ingestable_files(root: Path) -> List[Path]:
    files: List[Path] = []
    for p in root.rglob("*"):
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTS:
            files.append(p)
    return sorted(files)


def main():
    print("Step 0: Pre-flight")
    required = {
        "AZURE_OPENAI_ENDPOINT": AZURE_OPENAI_ENDPOINT,
        "AZURE_OPENAI_API_KEY": AZURE_OPENAI_API_KEY,
        "AZURE_OPENAI_API_VERSION": AZURE_OPENAI_API_VERSION,
        "AZURE_CHAT_API_VERSION": AZURE_CHAT_API_VERSION,
        "PG_CONN": PG_CONN,
    }
    for name, val in required.items():
        if _is_missing(val):
            print(f"ERROR: Please set a real value for {name} (current type: {type(val).__name__}).")
            sys.exit(1)

    ensure_registry_table()

    print(f"Step 1: Scanning folder: {DOC_ROOT}")
    files = list_ingestable_files(DOC_ROOT)
    if not files:
        print("No supported files found. Exiting.")
        return
    print(f"Found {len(files)} file(s).")

    print("Step 2: Init embeddings and vector store")
    embeddings = build_embeddings_client()
    vectorstore = PGVector(
        connection=PG_CONN,
        embeddings=embeddings,
        collection_name=COLLECTION_NAME,
        use_jsonb=True
    )

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP
    )

    total_added = 0
    total_skipped = 0
    total_replaced = 0

    to_add_docs: List[Document] = []
    to_add_ids: List[str] = []

    for f in tqdm(files, desc="Indexing files", unit="file"):
        # Extract & hash
        try:
            text = extract_text_from_any(f)
        except Exception as e:
            print(f"\nWARN: Failed to read {f}: {e}")
            continue

        if not text.strip():
            print(f"\nNOTE: Empty text after extraction: {f}")
            total_skipped += 1
            continue

        mtime = int(f.stat().st_mtime)
        chash = hash_text(text)
        src = str(f.resolve())
        old_hash = get_existing_hash(src)

        doc_id = f"{DOC_ROOT.as_posix()}::{f.relative_to(DOC_ROOT).as_posix()}"

        if old_hash and old_hash == chash:
            total_skipped += 1
            continue

        # If changed/new: try to delete old chunks for this doc_id
        try:
            # Some versions of langchain_postgres support 'filter', others 'where'
            try:
                vectorstore.delete(filter={"doc_id": doc_id})
            except TypeError:
                vectorstore.delete(where={"doc_id": doc_id})
            if old_hash:
                total_replaced += 1
        except Exception:
            # Ignore if delete not supported; we'll just add new chunks
            pass

        chunks = splitter.split_text(text)
        if not chunks:
            print(f"\nNOTE: No chunks produced for: {f}")
            continue

        for idx, chunk in enumerate(chunks):
            meta: Dict[str, str] = {
                "source_path": src,
                "doc_id": doc_id,
                "chunk_index": idx,
                "content_hash": chash,
                "file_mtime": mtime,
                "rel_path": str(f.relative_to(DOC_ROOT)),
                "file_ext": f.suffix.lower(),
                "ingested_at": datetime.utcnow().isoformat() + "Z",
            }
            to_add_docs.append(Document(page_content=chunk, metadata=meta))
            to_add_ids.append(f"{doc_id}#chunk-{idx}")

        if len(to_add_docs) >= BATCH_SIZE:
            vectorstore.add_documents(to_add_docs, ids=to_add_ids)
            total_added += len(to_add_docs)
            to_add_docs.clear()
            to_add_ids.clear()

        upsert_registry(src, chash, mtime, num_chunks=len(chunks))

    if to_add_docs:
        vectorstore.add_documents(to_add_docs, ids=to_add_ids)
        total_added += len(to_add_docs)

    print(f"Step 3: Ingest complete | added chunks: {total_added} | replaced docs: {total_replaced} | skipped unchanged: {total_skipped}")

    # Optional: quick RAG sanity check
    print("Step 4: Building retriever + QA chain")
    llm = build_chat_llm()
    prompt = ChatPromptTemplate.from_messages([
        ("system",
         "You are an assistant for question-answering tasks. "
         "Use the following pieces of retrieved context to answer the question. "
         "If you don't know the answer, say that you don't know. "
         "Use three sentences maximum and keep the answer concise.\n\n{context}"),
        ("human", "{input}"),
    ])
    qa_chain = create_stuff_documents_chain(llm, prompt)
    retriever = PGVector(
        connection=PG_CONN,
        embeddings=build_embeddings_client(),
        collection_name=COLLECTION_NAME,
        use_jsonb=True
    ).as_retriever(search_kwargs={"k": 4})
    rag_chain = create_retrieval_chain(retriever, qa_chain)

    test_q = "What is the main topic in the corpus under C:\\pp?"
    print("Step 5: Sample query:", test_q)
    resp = rag_chain.invoke({"input": test_q})
    print("Answer:", resp.get("answer", ""))

    print("Sources (retrieved chunks):")
    for i, d in enumerate(resp.get("context", [])[:5], 1):
        meta = d.metadata or {}
        head = d.page_content[:120].replace("\n", " ")
        print(f"  {i}. {meta.get('rel_path')} [chunk {meta.get('chunk_index')}] - {head}...")

if __name__ == "__main__":
    main()
